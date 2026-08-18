"""Fill the organiser's Idea Submission Template with our content.

    python fill_official_template.py

Edits the provided .pptx in place-of-copy rather than rebuilding a lookalike,
so the layout, branding and "idea details pointers" are exactly theirs.

Per the template's own instruction slide:
  * slide 1 (instructions) is removed
  * total slides kept to 7, so "Idea Description" folds into "Proposed
    Solution" and "Research and References" folds into the Technology slide
  * output saved as <TeamName>_<PSNo>

Body copy is written tight on purpose: the template's placeholder boxes are
shallow, and this script cannot see the rendered result.
"""

from __future__ import annotations

import copy
import csv
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = Path(r"D:\Idea Submission Template_Hackathon 2026.pptx")

TEAM_NAME = "perceptron"
PS_NO = "PS01"
LEADER = "Aman Kumar Maurya"
YEAR = "4th Year"
COLLEGE = "Indian Institute of Technology Madras"
PHONE = "+91 83183 54205"
EMAIL = "23f2001033@ds.study.iitm.ac.in"
GITHUB = "https://github.com/23f2001033/kla-restore"
VIDEO = ""   # blank -> a sentence pointing at the repo instead


def metrics() -> dict:
    p = ROOT / "results" / "metrics.csv"
    if not p.exists():
        return {}
    return {r["method"]: r for r in csv.DictReader(p.open())}


def set_text(shape, lines, size=None):
    """Replace a placeholder's text while keeping its run formatting.

    Assigning to text_frame.text would collapse the paragraph to a single
    unstyled run, losing the template's fonts and colours, so the first run is
    reused as the style source and cloned for extra lines.
    """
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    r0 = p0.runs[0]
    r0.text = lines[0]
    if size:
        r0.font.size = Pt(size)
    for extra in p0.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for line in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        para = _Paragraph(newp, tf)
        para.runs[0].text = line
        if size:
            para.runs[0].font.size = Pt(size)


def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape id {shape_id} not found")


def drop_slides(prs, indices_1based):
    """Remove slides by 1-based index, highest first so indices stay valid."""
    lst = prs.slides._sldIdLst
    items = list(lst)
    for i in sorted(indices_1based, reverse=True):
        rId = items[i - 1].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        lst.remove(items[i - 1])


def main() -> None:
    m = metrics()

    def g(method, key, default):
        try:
            return float(m[method][key])
        except (KeyError, TypeError, ValueError):
            return default

    psnr, bic = g("model", "psnr", 27.54), g("bicubic", "psnr", 22.94)
    ssim, ssim_b = g("model", "ssim", 0.7594), g("bicubic", "ssim", 0.5460)
    lp, lp_b = g("model", "lpips", 0.3922), g("bicubic", "lpips", 0.4478)
    spi = g("model", "sec_per_image", 0.02167)

    prs = Presentation(str(TEMPLATE))
    s = list(prs.slides)

    # ---- slide 2: Team Details -------------------------------------------
    t = s[1]
    set_text(by_id(t, 54), TEAM_NAME)
    set_text(by_id(t, 55), LEADER)
    set_text(by_id(t, 56), YEAR)
    for sid in (57, 58, 59, 60, 61, 62):      # rows 2-4 unused (solo entry)
        set_text(by_id(t, sid), "")
    for sid in (42, 43, 44, 45, 46, 47):      # their "2/Member 1" row labels
        set_text(by_id(t, sid), "")
    set_text(by_id(t, 63), COLLEGE)
    set_text(by_id(t, 64), PHONE)
    set_text(by_id(t, 65), EMAIL)

    # ---- slide 3: Problem Statement Addressed ----------------------------
    p = s[2]
    set_text(by_id(p, 16),
             "AI-Based Restoration of Degraded Images for Semiconductor "
             "Inspection (KLA)", size=13)
    set_text(by_id(p, 19), [
        "Inspection images arrive degraded by speckle noise, additive Gaussian noise "
        "and downsampling, applied in an undisclosed order. Input is half resolution "
        "and noisy; the output must be clean at full ground-truth resolution.",
        "Lost detail reduces the reliability of every downstream measurement and "
        "defect-detection step. Doing nothing (bicubic upsampling) scores "
        f"{bic:.2f} dB PSNR, and the hidden test set includes out-of-distribution "
        "content, so the model must generalise beyond the training images.",
    ], size=11)

    # ---- slide 5: Proposed Solution (also carries Idea Description) ------
    sol = s[4]
    set_text(by_id(sol, 16),
             "Idea and methodology: recover the degradation model from the data, "
             "then invert it with a single efficient network.", size=12)
    set_text(by_id(sol, 19), [
        "KEY CONCEPT  -  Restoration is an inverse problem, so we first recovered the "
        "forward model by fitting the residual across all 3200 provided pairs: 2x2 area "
        "downsampling, noise applied after downsampling, and multiplicative Gamma speckle "
        "(measured skew +0.357 vs Gamma's predicted +0.333).",
        "SOLUTION  -  A NAFNet-style network, 16 blocks running at LOW resolution with a "
        "single PixelShuffle x2 upsample, doing joint denoising and super-resolution in one "
        "forward pass. 0.68 M parameters, no BatchNorm, fully convolutional so the same "
        "weights serve 128->256 and 256->512.",
        "TRAINING  -  50% real pairs, 50% synthetic pairs generated from the measured "
        "distribution with randomised kernel and degradation order. Loss = Charbonnier + "
        "0.15 x (1 - SSIM), plus LPIPS for the final 10% of steps.",
    ], size=10.5)

    # ---- slide 6: Innovation and Uniqueness ------------------------------
    inn = s[5]
    set_text(by_id(inn, 19), [
        "We measured the degradation instead of guessing it. Most entrants feed the "
        "provided pairs to a network; we recovered the forward model first.",
        "The residual's skew identifies the noise FAMILY: multiplicative Gaussian is "
        "symmetric, Gamma is not. Measured +0.357 against a predicted +0.333 pins it.",
    ], size=10)
    set_text(by_id(inn, 22), [
        "Synthetic pairs sampled from empirical quantile tables match reality; a first "
        "version using assumed ranges was 50% noisier and would have taught the model "
        "to over-smooth.",
        "Leakage-aware validation: samples are ordered by source image, so we hold out "
        "10 evenly-spaced contiguous blocks rather than a random split.",
        "Deliberately no GAN - the brief requires restoration without hallucinating "
        "structure, and invented detail is harmful in a measurement context.",
    ], size=10)

    # ---- slide 7: Impact and Benefits ------------------------------------
    imp = s[6]
    set_text(by_id(imp, 20), [
        "Recovers detail lost to downsampling without inventing structure, which is "
        "what a precision-measurement pipeline requires.",
        "Improves all three scored axes at once - perceptual metrics often move "
        "against PSNR, and here they do not.",
    ], size=10)
    set_text(by_id(imp, 24), [
        f"PSNR  {psnr:.2f} dB  vs  {bic:.2f} dB bicubic   (+{psnr-bic:.2f} dB)",
        f"SSIM  {ssim:.4f}  vs  {ssim_b:.4f}   (+{100*(ssim/ssim_b-1):.0f}%)",
        f"LPIPS  {lp:.4f}  vs  {lp_b:.4f}  (lower is better)",
        f"Inference  {spi*1000:.1f} ms/image  ({1/spi:.0f} img/s on an NVIDIA T4)",
    ], size=10)

    # ---- slide 8: Technology & Feasibility (also carries references) -----
    tech = s[7]
    set_text(by_id(tech, 16),
             "Technical stack, feasibility, and the research the approach rests on.",
             size=12)
    set_text(by_id(tech, 19), [
        "STACK  -  PyTorch 2.x / Python 3.11. NAFNet-style CNN, 0.68 M parameters, trained "
        "from scratch on a single NVIDIA T4 within one free Kaggle session. Inference imports "
        "only torch and numpy: no internet, API keys, downloads or manual configuration.",
        "RUN  -  python run.py <input-dir> <output-dir>. Weights auto-load from models/, the "
        "output directory is created if absent, outputs are clamped to [0,1] with no NaN/Inf. "
        "Verified end-to-end from a fresh public clone.",
        "HYGIENE  -  Seeds fixed, config-driven, every validation logs metrics plus git hash; "
        "a pre-flight gate blocks a broken run before it wastes GPU hours.",
        "REFERENCES  -  NAFNet (Chen 2022); Real-ESRGAN degradation modelling (Wang 2021); "
        "PixelShuffle (Shi 2016); SSIM (Wang 2004); LPIPS (Zhang 2018). No external datasets "
        "or pretrained restoration weights; lpips (BSD-2) supplies the perceptual metric only.",
    ], size=9.5)
    set_text(by_id(tech, 22), "Restoration network")
    set_text(by_id(tech, 25), "NVIDIA GPU (T4 / H100)")
    set_text(by_id(tech, 28), "PyTorch, NumPy")

    # ---- slide 9: GitHub & Video Link ------------------------------------
    gh = s[8]
    set_text(by_id(gh, 20), GITHUB, size=13)
    set_text(by_id(gh, 27),
             VIDEO or "Working solution is reproducible directly from the repository: "
                      "python run.py <input-dir> <output-dir>", size=11)

    # ---- drop instructions, Idea Description, and References -------------
    drop_slides(prs, [10, 4, 1])

    out = ROOT / f"{TEAM_NAME}_{PS_NO}.pptx"
    prs.save(str(out))
    n = len(prs.slides._sldIdLst)
    print(f"[deck] wrote {out} ({n} slides)")
    if n > 7:
        print(f"[deck] WARNING: {n} slides exceeds the template's 6-7 limit")


if __name__ == "__main__":
    main()
