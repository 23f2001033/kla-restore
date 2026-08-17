"""Generate the KLA solution presentation from the recorded results.

    python make_presentation.py

Reads results/metrics.csv, results/run_log.csv and results/qualitative.png when
they exist, so the deck regenerates with real numbers the moment training
finishes -- no hand-editing of slides. Any metric that has not been produced yet
renders as "pending" rather than a fabricated value.

Follows the 12-slide structure prescribed in KLA's participant help document.
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- palette: instrumentation / signal, not generic corporate blue -----------
NAVY = RGBColor(0x0B, 0x1B, 0x2B)      # dark ground for title + section slides
LIGHT = RGBColor(0xF7, 0xF9, 0xFB)     # light ground for content slides
CYAN = RGBColor(0x00, 0xB4, 0xD8)      # accent: restored signal
AMBER = RGBColor(0xFF, 0xB7, 0x03)     # accent: degradation / warnings
INK = RGBColor(0x1A, 0x27, 0x33)       # body text on light
MUTED = RGBColor(0x5A, 0x6B, 0x7A)     # captions
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD8, 0xE1, 0xE8)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

W, H = 13.333, 7.5
MARGIN = 0.7

ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / "results"


# --------------------------------------------------------------------------- #
# data loading
# --------------------------------------------------------------------------- #
def load_metrics() -> dict[str, dict]:
    path = RESULTS / "metrics.csv"
    if not path.exists():
        return {}
    out = {}
    with path.open() as f:
        for row in csv.DictReader(f):
            out[row["method"]] = row
    return out


def load_run() -> dict:
    path = RESULTS / "run_log.csv"
    if not path.exists():
        return {}
    rows = list(csv.DictReader(path.open()))
    if not rows:
        return {}
    best = max(rows, key=lambda r: float(r["val_psnr"] or -1))
    return {"rows": rows, "best": best, "last": rows[-1]}


def fmt(value, spec: str = "{:.2f}", pending: str = "pending") -> str:
    try:
        f = float(value)
    except (TypeError, ValueError):
        return pending
    if f != f:  # NaN
        return pending
    return spec.format(f)


# --------------------------------------------------------------------------- #
# drawing helpers
# --------------------------------------------------------------------------- #
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size=14, bold=False, color=INK, font=BODY_FONT,
         space_after=6, first=False, align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return p


def bullets(tf, items, size=13, color=INK, space_after=7, first=False):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = "•   " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = BODY_FONT


def title(slide, text, sub=None, dark=False):
    tf = textbox(slide, MARGIN, 0.45, W - 2 * MARGIN, 0.9)
    para(tf, text, size=32, bold=True, color=WHITE if dark else INK,
         font=HEAD_FONT, first=True, space_after=2)
    if sub:
        tf2 = textbox(slide, MARGIN, 1.28, W - 2 * MARGIN, 0.4)
        para(tf2, sub, size=13, color=CYAN if dark else MUTED, first=True)


def card(slide, x, y, w, h, fill=CARD, line=RULE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except (IndexError, AttributeError):
            pass
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    return shape, tf


def stat(slide, x, y, w, value, label, color=CYAN, vsize=30, note=None):
    tf = textbox(slide, x, y, w, 1.15, align=PP_ALIGN.LEFT)
    para(tf, value, size=vsize, bold=True, color=color, font=HEAD_FONT,
         first=True, space_after=1)
    para(tf, label, size=10.5, color=MUTED, space_after=0)
    if note:
        para(tf, note, size=9.5, color=MUTED, italic=True, space_after=0)


def table(slide, x, y, w, rows, col_w, head_bg=NAVY, size=11, row_h=0.34):
    """Lightweight table: header row + body rows, drawn as text boxes."""
    hy = y
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(hy),
                                 Inches(w), Inches(row_h))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = head_bg
    hdr.line.fill.background()
    hdr.shadow.inherit = False
    hdr.text_frame.text = ""

    cx = x
    for j, cell in enumerate(rows[0]):
        tf = textbox(slide, cx + 0.12, hy + 0.055, col_w[j] - 0.16, row_h,
                     anchor=MSO_ANCHOR.MIDDLE)
        para(tf, str(cell), size=size, bold=True, color=WHITE, first=True, space_after=0)
        cx += col_w[j]

    for i, row in enumerate(rows[1:]):
        ry = hy + row_h * (i + 1)
        if i % 2 == 0:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(ry),
                                          Inches(w), Inches(row_h))
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(0xEF, 0xF3, 0xF6)
            band.line.fill.background()
            band.shadow.inherit = False
            band.text_frame.text = ""
        cx = x
        for j, cell in enumerate(row):
            emph = isinstance(cell, str) and cell.startswith("*")
            txt = cell[1:] if emph else str(cell)
            tf = textbox(slide, cx + 0.12, ry + 0.055, col_w[j] - 0.16, row_h,
                         anchor=MSO_ANCHOR.MIDDLE)
            para(tf, txt, size=size, bold=emph, color=CYAN if emph else INK,
                 first=True, space_after=0)
            cx += col_w[j]
    return hy + row_h * len(rows)


def footer(slide, text, dark=False):
    tf = textbox(slide, MARGIN, H - 0.52, W - 2 * MARGIN, 0.3)
    para(tf, text, size=9, color=MUTED if not dark else RGBColor(0x7A, 0x8B, 0x9A),
         italic=True, first=True, space_after=0)


# --------------------------------------------------------------------------- #
# slides
# --------------------------------------------------------------------------- #
def slide_title(prs, m):
    s = blank(prs)
    bg(s, NAVY)
    tf = textbox(s, MARGIN, 2.05, 9.6, 1.6)
    para(tf, "Restoring Degraded", size=42, bold=True, color=WHITE,
         font=HEAD_FONT, first=True, space_after=2)
    para(tf, "Inspection Images", size=42, bold=True, color=CYAN,
         font=HEAD_FONT, space_after=10)
    tf2 = textbox(s, MARGIN, 3.95, 9.8, 1.0)
    para(tf2, "Joint speckle/Gaussian denoising and 2x super-resolution, "
              "built on a degradation model recovered from the data itself.",
         size=15, color=RGBColor(0xC9, 0xD6, 0xE0), first=True)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(3.72),
                              Inches(1.5), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()
    line.shadow.inherit = False

    tf3 = textbox(s, MARGIN, 5.5, 11.5, 0.9)
    para(tf3, "KLA Problem Statement  |  SEMICON India Hackathon 2026",
         size=12, bold=True, color=WHITE, first=True, space_after=3)
    para(tf3, "Team submission  ·  github.com/23f2001033/kla-restore",
         size=11, color=MUTED)


def slide_problem(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "The restoration task",
          "Three degradations, applied in an undisclosed order")

    _, tf = card(s, MARGIN, 1.85, 6.0, 3.3)
    para(tf, "What we are given", size=15, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=9)
    bullets(tf, [
        "3200 paired samples: clean GT and its degraded counterpart",
        "GT 256×256 float32 in [0,1]; NoisyLR 128×128",
        "Degradations: speckle noise, additive Gaussian noise, downsampling",
        "Order undisclosed — the model must not depend on knowing it",
        "NoisyLR deliberately exceeds [0,1] (observed [-0.28, 2.16])",
    ], size=12.5)

    _, tf2 = card(s, MARGIN + 6.35, 1.85, 5.6, 3.3)
    para(tf2, "What is scored", size=15, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=9)
    bullets(tf2, [
        "Restoration quality: PSNR, SSIM and LPIPS on hidden ground truth",
        "End-to-end throughput on an NVIDIA H100 — including script "
        "startup, disk I/O and pre/post-processing",
        "Training and compute hygiene: reproducibility, clean experiments",
        "Hidden test set spans in-distribution and out-of-distribution content",
    ], size=12.5)

    stat(s, MARGIN, 5.45, 3.0, "23.18 dB", "bicubic upsample of NoisyLR",
         color=AMBER, note="the do-nothing baseline")
    stat(s, MARGIN + 3.3, 5.45, 3.4, "32.25 dB", "perfect denoise + bicubic",
         color=MUTED, note="interpolation-only ceiling")
    stat(s, MARGIN + 7.0, 5.45, 4.4, "> 32.25 dB", "target for a model that "
         "truly learns SR", color=CYAN, note="must beat interpolation, not match it")
    footer(s, "Baselines measured on the held-out split; reproduce with "
              "python analyze_degradation.py")


def slide_degradation(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "We recovered the degradation model from the data",
          "Rather than guessing at it — every number below is reproducible")

    rows = [
        ["Property", "Finding", "Evidence"],
        ["Downsampling", "*2×2 area (box) average",
         "residual std 0.0910 vs 0.0924 bicubic, 0.1006 nearest"],
        ["Noise placement", "*applied after downsampling",
         "residual is spatially white (lag-1 autocorr −0.05)"],
        ["Speckle family", "*multiplicative Gamma",
         "skew of r/d = +0.357 vs Gamma's predicted 2/√36 = +0.333"],
        ["Speckle strength", "L: p5 22, median 36, p95 57",
         "per-image fit of var(r) = d²/L + σ²"],
        ["Gaussian strength", "σ: median 0.021, p95 0.073",
         "21.3% of images below 0.005"],
    ]
    table(s, MARGIN, 1.95, W - 2 * MARGIN, rows, [2.3, 3.5, 6.13], row_h=0.42)

    _, tf = card(s, MARGIN, 4.45, 11.93, 1.55, fill=RGBColor(0xE8, 0xF7, 0xFB),
                 line=CYAN)
    para(tf, "Why this mattered", size=13.5, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=6)
    para(tf, "Multiplicative Gaussian speckle would be symmetric; the measured "
             "+0.357 skew is what pins the family to Gamma. Fixing the forward "
             "model is what let us generate synthetic pairs that match the real "
             "distribution instead of a plausible-looking guess — the first "
             "attempt was 50% too noisy (residual std 0.137 vs 0.091) and would "
             "have biased the model toward over-smoothing.",
         size=12, color=INK)
    footer(s, "Full derivation in results/degradation_analysis.md")


def slide_pipeline(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "End-to-end pipeline", "One model, one forward pass, no staging")

    steps = [
        ("Pack", "6400 .npy files → two contiguous memmaps",
         "removes per-file I/O from every epoch"),
        ("Degrade", "50% real pairs / 50% synthetic",
         "randomised kernel, order, L and σ"),
        ("Train", "NAF body at LR res + PixelShuffle ×2",
         "Charbonnier + SSIM (+LPIPS fine-tune)"),
        ("Infer", "batched, fp16, threaded I/O",
         "clamped to [0,1] before writing"),
    ]
    x = MARGIN
    cw = 2.82
    for i, (h, sub, note) in enumerate(steps):
        _, tf = card(s, x, 2.15, cw, 2.15)
        para(tf, f"0{i+1}", size=11, bold=True, color=CYAN, first=True, space_after=4)
        para(tf, h, size=17, bold=True, color=INK, font=HEAD_FONT, space_after=7)
        para(tf, sub, size=11.5, color=INK, space_after=6)
        para(tf, note, size=10, color=MUTED, italic=True, space_after=0)
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + cw + 0.04),
                                    Inches(3.03), Inches(0.28), Inches(0.36))
            ar.fill.solid()
            ar.fill.fore_color.rgb = CYAN
            ar.line.fill.background()
            ar.shadow.inherit = False
            ar.text_frame.text = ""
        x += cw + 0.36

    _, tf = card(s, MARGIN, 4.75, 11.93, 1.5, fill=NAVY, line=None)
    para(tf, "Scale-agnostic by construction", size=13.5, bold=True,
         color=CYAN, font=HEAD_FONT, first=True, space_after=6)
    para(tf, "The network is fully convolutional with no fixed-size assumption, "
             "so the same weights handle 128→256 and 256→512. The brief states "
             "test data may contain either. Verified end-to-end: a 256×256 input "
             "emits 512×512 with no shape assertion.",
         size=12, color=RGBColor(0xD5, 0xE2, 0xEC))
    footer(s, "python inference.py --input_dir <dir> --output_dir <dir>")


def slide_augmentation(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Synthetic degradation engine",
          "Sampled from the measured distribution, then deliberately widened")

    _, tf = card(s, MARGIN, 1.9, 5.75, 2.5)
    para(tf, "How parameters are sampled", size=14, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=8)
    para(tf, "(L, σ) are drawn from empirical quantile tables fitted on all "
             "3200 real pairs, then multiplied by a random widening factor so "
             "the tails extend past what was observed. This reproduces the true "
             "distribution by construction rather than assuming a parametric form.",
         size=11.5, color=INK, space_after=7)
    para(tf, "Kernel ∈ {area, bilinear, bicubic, gaussian+stride};  order "
             "randomised across the three degradations.",
         size=11, color=MUTED, italic=True)

    rows = [
        ["Source", "residual std", "skew r/d", "per-image p10 / p50 / p90"],
        ["real pairs", "0.0852", "+0.367", "0.045 / 0.079 / 0.114"],
        ["synthetic (exact)", "0.0945", "+0.299", "0.050 / 0.081 / 0.135"],
        ["synthetic (widened)", "0.1010", "+0.483", "0.045 / 0.086 / 0.145"],
    ]
    table(s, MARGIN + 6.1, 1.9, 6.2, rows, [1.85, 1.15, 1.0, 2.2], size=10.5,
          row_h=0.4)

    _, tf2 = card(s, MARGIN + 6.1, 3.72, 6.2, 0.72, fill=RGBColor(0xE8, 0xF7, 0xFB),
                  line=CYAN)
    para(tf2, "Synthetic pairs match real statistics closely; the widened "
              "variant extends both tails for OOD robustness.",
         size=11, color=INK, first=True, space_after=0)

    _, tf3 = card(s, MARGIN, 4.7, 11.93, 1.5)
    para(tf3, "Training mix and augmentation", size=13.5, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=7)
    bullets(tf3, [
        "50% real provided pairs (anchors the exact test distribution) / "
        "50% synthetic (widens it)",
        "64×64 LR crops → 128×128 GT, random flips and 90° rotations; "
        "crop-based training is what makes the model scale-agnostic",
        "Synthetic LR is never clipped — the real data is not either, and "
        "clipping would create a train/test mismatch on exactly those values",
    ], size=11.5)
    footer(s, "src/degrade.py")


def slide_model(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Architecture", "NAFNet-style body at LR resolution, "
                             "single PixelShuffle ×2 head")

    _, tf = card(s, MARGIN, 1.9, 5.5, 3.4, fill=NAVY, line=None)
    para(tf, "input  (B, 1, H, W)", size=12, bold=True, color=CYAN,
         font=BODY_FONT, first=True, space_after=8)
    for line_txt in [
        "conv3×3  →  width 64",
        "16 × NAFBlock",
        "   depthwise conv + SimpleGate",
        "   simplified channel attention",
        "conv3×3  (long residual)",
        "conv3×3  →  PixelShuffle(2)",
        "conv3×3  →  1 channel",
        "+ bilinear(input, ×2)",
    ]:
        para(tf, line_txt, size=12, color=RGBColor(0xD5, 0xE2, 0xEC), space_after=5)
    para(tf, "output  (B, 1, 2H, 2W)", size=12, bold=True, color=CYAN, space_after=0)

    _, tf2 = card(s, MARGIN + 5.85, 1.9, 6.08, 3.4)
    para(tf2, "Design rationale", size=14, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=8)
    bullets(tf2, [
        "Body runs at LR resolution and upsamples once at the end — ~4× "
        "cheaper than a U-Net at output resolution, and throughput is scored",
        "NAF blocks reach near-transformer quality with pure convolution; "
        "SwinIR/Restormer gain little at 128×128 tiles and cost real H100 time",
        "No BatchNorm anywhere — channel-wise LayerNorm only, so the "
        "evaluator's batch size cannot change our outputs",
        "Per-block residual scales initialise at zero: a 16-block stack starts "
        "as identity and trains stably from scratch",
        "No GAN — the brief asks for restoration without hallucinating "
        "structure",
    ], size=11.5, space_after=6)

    stat(s, MARGIN, 5.6, 2.6, "0.68 M", "parameters", color=CYAN, vsize=26)
    stat(s, MARGIN + 3.0, 5.6, 3.2, "4.8e−07", "max |batch8 − batch1|",
         color=CYAN, vsize=26, note="batch-size invariance verified")
    stat(s, MARGIN + 6.9, 5.6, 5.0, "128→256 and 256→512",
         "same weights, no reconfiguration", color=CYAN, vsize=22)
    footer(s, "src/model.py")


def slide_loss(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Loss and training setup",
          "Pixel fidelity, structure and perception — weighted for what is scored")

    _, tf = card(s, MARGIN, 1.9, 5.75, 2.75)
    para(tf, "Objective", size=14, bold=True, color=INK, font=HEAD_FONT,
         first=True, space_after=8)
    para(tf, "L  =  Charbonnier  +  0.15 · (1 − SSIM)",
         size=13.5, bold=True, color=CYAN, space_after=4)
    para(tf, "+  0.05 · LPIPS   for the final 10% of steps",
         size=12, color=MUTED, space_after=9)
    bullets(tf, [
        "Charbonnier over L2: better PSNR/SSIM trade-off and robust to the "
        "heavy-tailed speckle outliers",
        "Single-scale SSIM, not MS-SSIM — MS-SSIM needs ≥161 px and the GT "
        "crops are 128×128; it would fail or silently degrade",
        "LPIPS is scored, but VGG every step costs ~30% throughput — a short "
        "fine-tune captures most of the benefit",
    ], size=11.5, space_after=6)

    rows = [
        ["Setting", "Value"],
        ["Optimiser", "AdamW, betas (0.9, 0.999), wd 1e-4"],
        ["Learning rate", "5e-4, 2000-step warmup, cosine to 1e-6"],
        ["Batch / precision", "32, AMP, channels_last"],
        ["Stability", "grad-norm clip 1.0, non-finite loss guard"],
        ["Checkpointing", "validate every 5000 steps, keep best by PSNR"],
    ]
    table(s, MARGIN + 6.1, 1.9, 6.2, rows, [2.1, 4.1], size=11, row_h=0.4)

    _, tf3 = card(s, MARGIN, 5.0, 11.93, 1.25, fill=RGBColor(0xFF, 0xF6, 0xE3),
                  line=AMBER)
    para(tf3, "Stability incident and fix", size=13, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=5)
    para(tf3, "An initial run diverged to a non-finite loss exactly as warmup "
              "ended and the learning rate peaked. Root cause: AdamW was "
              "configured with betas=(0.9, 0.9); the fast second-moment decay "
              "made the per-parameter step size volatile enough that one "
              "gradient spike at peak LR corrupted the weights irrecoverably. "
              "Fixed to (0.9, 0.999), lowered the peak LR to 5e-4, and added a "
              "guard that skips non-finite batches and aborts after 20 "
              "consecutive ones rather than silently consuming the budget.",
         size=11, color=INK)
    footer(s, "src/losses.py, train.py, configs/base.yaml")


def slide_experiments(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Experiment tracking and validation design",
          "Reproducibility as a first-class deliverable")

    _, tf = card(s, MARGIN, 1.9, 5.75, 3.05)
    para(tf, "The validation split is group-aware", size=14, bold=True,
         color=INK, font=HEAD_FONT, first=True, space_after=8)
    para(tf, "KLA's own slides show sample 000000 → source 0001.png and "
             "sample 000500 → source 0186.png: samples are ordered by source "
             "image, roughly 2-3 per source.",
         size=11.5, color=INK, space_after=7)
    bullets(tf, [
        "A random per-index split therefore leaks near-duplicate content "
        "into training",
        "A contiguous tail slice hands validation sources seen nowhere else",
        "We hold out 10 evenly-spaced contiguous blocks of 32 — contiguous "
        "keeps same-source samples together, spacing covers all content types",
        "Split frozen in configs/base.yaml; train ∩ val = ∅ asserted at load",
    ], size=11.5, space_after=6)

    _, tf2 = card(s, MARGIN + 6.1, 1.9, 6.2, 3.05)
    para(tf2, "What every run records", size=14, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=8)
    bullets(tf2, [
        "run id, git hash, seed, step, LR, loss, all three metrics, parameter "
        "count and the full config as JSON — appended to results/run_log.csv "
        "at every validation",
        "Seeds fixed for Python, NumPy and Torch",
        "Every hyperparameter lives in configs/base.yaml; nothing hard-coded",
        "The checkpoint embeds its own model config, so inference rebuilds the "
        "architecture without reading any config file",
        "Schedule length is calibrated from measured throughput, so cosine "
        "reaches its minimum exactly when the budget ends",
    ], size=11.5, space_after=6)

    _, tf3 = card(s, MARGIN, 5.25, 11.93, 1.0, fill=NAVY, line=None)
    para(tf3, "Pre-flight gate", size=12.5, bold=True, color=CYAN,
         font=HEAD_FONT, first=True, space_after=5)
    para(tf3, "train.py --overfit 2 --steps 2000 must exceed 35 dB on two "
              "images before any long run starts, and exits non-zero on "
              "failure so it can gate a script. Crop alignment is verified "
              "independently: aligned residual 0.078 vs 0.128 for a 1-pixel "
              "shifted control.",
         size=11, color=RGBColor(0xD5, 0xE2, 0xEC))
    footer(s, "results/run_log.csv")


def slide_results(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Restoration quality", "Held-out validation split, 320 images "
                                    "never seen in training")

    model = m.get("model", {})
    bic = m.get("bicubic", {})
    has = bool(model)

    rows = [
        ["Method", "PSNR (dB)", "SSIM", "LPIPS"],
        ["bicubic upsample (baseline)",
         fmt(bic.get("psnr"), "{:.2f}") if bic else "23.18",
         fmt(bic.get("ssim"), "{:.4f}") if bic else "pending",
         fmt(bic.get("lpips"), "{:.4f}") if bic else "pending"],
        ["perfect denoise + bicubic (ceiling)", "32.25", "—", "—"],
        ["*this model",
         "*" + fmt(model.get("psnr"), "{:.2f}"),
         "*" + fmt(model.get("ssim"), "{:.4f}"),
         "*" + fmt(model.get("lpips"), "{:.4f}")],
    ]
    table(s, MARGIN, 2.0, 11.93, rows, [5.0, 2.4, 2.3, 2.23], size=12, row_h=0.44)

    if has:
        gain = float(model.get("psnr", 0)) - float(bic.get("psnr", 23.18) or 23.18)
        stat(s, MARGIN, 4.15, 3.4, f"+{gain:.2f} dB", "over the bicubic baseline",
             color=CYAN)
    else:
        stat(s, MARGIN, 4.15, 3.4, "pending", "final training run in progress",
             color=MUTED)

    ood_m, ood_b = m.get("ood_model"), m.get("ood_bicubic")
    if ood_m and ood_b:
        stat(s, MARGIN + 3.7, 4.15, 4.0,
             f"{float(ood_m['psnr']):.2f} dB", "out-of-distribution probe",
             color=CYAN, note=f"vs {float(ood_b['psnr']):.2f} dB bicubic")
    else:
        stat(s, MARGIN + 3.7, 4.15, 4.0, "OOD probe", "external images "
             "degraded with our engine", color=MUTED,
             note="proxy for the hidden OOD split")

    _, tf = card(s, MARGIN, 5.35, 11.93, 1.15, fill=RGBColor(0xE8, 0xF7, 0xFB),
                 line=CYAN)
    para(tf, "Metric conventions — stated because implementations differ",
         size=12.5, bold=True, color=INK, font=HEAD_FONT, first=True, space_after=5)
    para(tf, "PSNR: data_range=1.0 on the clipped prediction, averaged per image.   "
             "SSIM: 11×11 Gaussian window, σ=1.5, data_range=1.0 — equivalent "
             "to skimage's gaussian_weights=True; skimage's default 7×7 uniform "
             "window yields a different number.   LPIPS: LPIPS-VGG on "
             "3-channel-replicated grayscale, inputs mapped to [−1, 1].",
         size=10.5, color=INK)
    footer(s, "results/metrics.csv  ·  python evaluate.py --weights weights/best.pt")


def slide_runtime(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Throughput", "Scored end-to-end: startup, I/O and "
                           "pre/post-processing all count")

    _, tf = card(s, MARGIN, 1.9, 5.75, 3.4)
    para(tf, "What we optimised", size=14, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=8)
    bullets(tf, [
        "Only torch, numpy and argparse imported at module level — no yaml, "
        "lpips, skimage or matplotlib, each of which costs startup time for "
        "no benefit",
        "Model config travels inside the checkpoint, so there is no config parse",
        "Files read and written on a thread pool (npy I/O releases the GIL)",
        "Inputs grouped by resolution and run in batches under inference_mode, "
        "with channels_last, pinned memory, non-blocking transfers and fp16",
        "torch.compile off by default — 30-60 s of compile will not amortise "
        "over a small test set; available behind --compile",
        "×8 self-ensemble available behind --tta: ~+0.2 dB for ~8× the "
        "runtime, off by default since throughput is scored",
    ], size=11, space_after=6)

    _, tf2 = card(s, MARGIN + 6.1, 1.9, 6.2, 3.4)
    para(tf2, "Measured", size=14, bold=True, color=INK, font=HEAD_FONT,
         first=True, space_after=8)
    mm = m.get("model", {})
    spi = mm.get("sec_per_image")
    if spi:
        try:
            v = float(spi)
            para(tf2, f"{v*1000:.1f} ms", size=30, bold=True, color=CYAN,
                 font=HEAD_FONT, space_after=2)
            para(tf2, "per image, model forward", size=11, color=MUTED, space_after=10)
            para(tf2, f"{1.0/v:.0f} images / second", size=13, bold=True,
                 color=INK, space_after=10)
        except (TypeError, ValueError, ZeroDivisionError):
            pass
    else:
        para(tf2, "pending", size=26, bold=True, color=MUTED, font=HEAD_FONT,
             space_after=10)
    para(tf2, "Full-pipeline wall clock is measured with the benchmark cell in "
              "the training notebook: it stages the validation split as loose "
              ".npy files and times inference.py from process start to the last "
              "file written — the same boundary KLA uses.",
         size=11, color=INK, space_after=8)
    para(tf2, "Hardware, batch size, driver, torch version and timing method "
              "are recorded alongside the number.",
         size=10.5, color=MUTED, italic=True)

    footer(s, "time python inference.py --input_dir <dir> --output_dir <dir>")


def slide_visuals(prs, m):
    s = blank(prs)
    bg(s, LIGHT)
    title(s, "Visual results and failure analysis",
          "Best and worst cases from the held-out split")

    fig = RESULTS / "qualitative.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(MARGIN), Inches(1.85),
                             height=Inches(4.0))
        tx = MARGIN + 5.4
        _, tf = card(s, tx, 1.85, W - tx - MARGIN, 4.0)
    else:
        _, tf = card(s, MARGIN, 1.9, 11.93, 3.9)

    para(tf, "Dominant failure mode", size=14, bold=True, color=INK,
         font=HEAD_FONT, first=True, space_after=8)
    para(tf, "Heavy speckle (low L) over fine high-frequency texture. Where the "
             "noise variance approaches the local signal variance, the model "
             "cannot separate texture from speckle and smooths both — the "
             "conservative choice, and the right one given the brief asks for "
             "restoration without hallucinating structure.",
         size=11.5, color=INK, space_after=9)
    para(tf, "Limitations", size=14, bold=True, color=INK, font=HEAD_FONT,
         space_after=7)
    bullets(tf, [
        "Trained on the released 256→128 pairs; the 512→256 case is handled "
        "by the fully-convolutional design and verified for shape correctness, "
        "but has no paired data to validate accuracy against",
        "The degradation engine assumes only the three named mechanisms",
        "No GAN or diffusion prior, so genuinely destroyed high-frequency "
        "detail is not invented — a deliberate accuracy/realism trade",
    ], size=11)
    footer(s, "results/qualitative.png — worst 4 and best 4 by PSNR")


def slide_conclusion(prs, m):
    s = blank(prs)
    bg(s, NAVY)
    title(s, "Summary", dark=True)

    items = [
        ("Measured, not assumed",
         "The degradation model was recovered from the data and validated "
         "statistically; synthetic pairs reproduce the real distribution."),
        ("Efficient by design",
         "0.68 M parameters, body at LR resolution, single upsample — quality "
         "without paying H100 throughput for it."),
        ("Reproducible end to end",
         "Config-driven, seeded, logged per validation; inference runs from a "
         "clean environment with two directory arguments."),
    ]
    x = MARGIN
    for h, body in items:
        _, tf = card(s, x, 1.95, 3.75, 2.0, fill=RGBColor(0x14, 0x2A, 0x3E),
                     line=RGBColor(0x24, 0x3E, 0x54))
        para(tf, h, size=14, bold=True, color=CYAN, font=HEAD_FONT,
             first=True, space_after=7)
        para(tf, body, size=11.5, color=RGBColor(0xD5, 0xE2, 0xEC))
        x += 3.99

    _, tf = card(s, MARGIN, 4.25, 11.93, 1.35, fill=RGBColor(0x14, 0x2A, 0x3E),
                 line=RGBColor(0x24, 0x3E, 0x54))
    para(tf, "External resources and licensing", size=13, bold=True, color=CYAN,
         font=HEAD_FONT, first=True, space_after=6)
    para(tf, "No external datasets or pretrained weights are used for the "
             "restoration model itself — it is trained from scratch on the "
             "official KLA data plus synthetic pairs derived from that same "
             "data. lpips (BSD-2-Clause) supplies the VGG perceptual loss for "
             "the final fine-tune and for reporting; its VGG16 backbone carries "
             "ImageNet-pretrained weights shipped with the package. Training "
             "degrades gracefully to Charbonnier + SSIM if it is unavailable, "
             "and inference.py never imports it.",
         size=11, color=RGBColor(0xD5, 0xE2, 0xEC))

    tf2 = textbox(s, MARGIN, 5.95, 11.93, 0.6)
    para(tf2, "github.com/23f2001033/kla-restore", size=15, bold=True,
         color=WHITE, font=HEAD_FONT, first=True, space_after=3)
    para(tf2, "README.md · train.py · inference.py · evaluate.py · "
              "analyze_degradation.py · configs/ · src/ · weights/ · results/",
         size=10.5, color=MUTED)


def main() -> None:
    m = load_metrics()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for fn in (slide_title, slide_problem, slide_degradation, slide_pipeline,
               slide_augmentation, slide_model, slide_loss, slide_experiments,
               slide_results, slide_runtime, slide_visuals, slide_conclusion):
        fn(prs, m)

    out = ROOT / "solution_presentation.pptx"
    prs.save(str(out))
    status = "with final metrics" if m.get("model") else "metrics pending"
    print(f"[ppt] wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, {status})")
    if not m.get("model"):
        print("[ppt] re-run after evaluate.py to populate the results slides")


if __name__ == "__main__":
    main()
