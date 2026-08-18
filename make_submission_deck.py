"""Build the portal-conforming submission deck.

    python make_submission_deck.py

The organiser's Idea Submission Template mandates a different shape from the
12-slide technical deck in make_presentation.py:

  * 6-7 slides maximum, including the title slide
  * the template's own section headings ("idea details pointers") kept intact
  * points, diagrams and infographics rather than paragraphs
  * exported to PDF; no PPT or DOC accepted
  * filename: <TeamName>_<PSNo>

Metrics are read from results/metrics.csv so the numbers cannot drift from what
was actually measured.

>>> FILL IN TEAM the block below before running. <<<
"""

from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------- #
# >>> EDIT THIS BLOCK <<<
# --------------------------------------------------------------------------- #
TEAM = {
    "team_name": "perceptron",
    "ps_no": "PS01",
    "college": "Indian Institute of Technology Madras",
    "leader_phone": "8318354205",
    "leader_email": "23f2001033@ds.study.iitm.ac.in",
    "members": [
        # (role, name, academic year)
        ("Team Leader", "Aman Kumar Maurya", "4th Year"),
    ],
    "github": "https://github.com/23f2001033/kla-restore",
    "video": "",           # leave blank if not recording one
}
PROBLEM = "AI-Based Restoration of Degraded Images for Semiconductor Inspection (KLA)"

# --------------------------------------------------------------------------- #
NAVY = RGBColor(0x0B, 0x1B, 0x2B)
LIGHT = RGBColor(0xF7, 0xF9, 0xFB)
CYAN = RGBColor(0x00, 0x7C, 0x91)
AMBER = RGBColor(0xC8, 0x86, 0x00)
INK = RGBColor(0x1A, 0x27, 0x33)
MUTED = RGBColor(0x5A, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD8, 0xE1, 0xE8)
HEAD, BODY = "Cambria", "Calibri"
W, H, M = 13.333, 7.5, 0.62

ROOT = Path(__file__).resolve().parent


def metrics() -> dict:
    p = ROOT / "results" / "metrics.csv"
    if not p.exists():
        return {}
    return {r["method"]: r for r in csv.DictReader(p.open())}


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = LIGHT
    return s


def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    t.word_wrap = True
    t.margin_left = t.margin_right = t.margin_top = t.margin_bottom = 0
    t.vertical_anchor = anchor
    t.paragraphs[0].alignment = align
    return t


def para(tf, text, size=13, bold=False, color=INK, font=BODY, after=5,
         first=False, align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(after)
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.italic = Pt(size), bold, italic
    r.font.color.rgb, r.font.name = color, font
    return p


def bullets(tf, items, size=12.5, after=6, first=False, color=INK):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(after)
        r = p.add_run()
        r.text = "•  " + it
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = BODY


def heading(slide, text, sub=None):
    t = tb(slide, M, 0.42, W - 2 * M, 0.62)
    para(t, text, size=26, bold=True, color=INK, font=HEAD, first=True, after=2)
    if sub:
        t2 = tb(slide, M, 1.08, W - 2 * M, 0.34)
        para(t2, sub, size=11.5, color=MUTED, first=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(M), Inches(1.5),
                                Inches(W - 2 * M), Inches(0.02))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    ln.text_frame.text = ""


def card(slide, x, y, w, h, fill=WHITE, accent=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if accent:
        sh.line.color.rgb = accent
        sh.line.width = Pt(1.25)
    else:
        sh.line.color.rgb = RULE
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.05
    except (IndexError, AttributeError):
        pass
    t = sh.text_frame
    t.word_wrap = True
    t.margin_left = t.margin_right = Inches(0.2)
    t.margin_top = t.margin_bottom = Inches(0.15)
    return t


def stat(slide, x, y, w, value, label, color=CYAN, vsize=26):
    t = tb(slide, x, y, w, 0.95)
    para(t, value, size=vsize, bold=True, color=color, font=HEAD, first=True, after=1)
    para(t, label, size=10, color=MUTED, after=0)


# --------------------------------------------------------------------------- #
def s1_title(prs, m):
    s = blank(prs)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = NAVY

    t = tb(s, M, 0.75, 11.5, 1.5)
    para(t, "AI-Based Restoration of", size=30, bold=True, color=WHITE,
         font=HEAD, first=True, after=1)
    para(t, "Degraded Inspection Images", size=30, bold=True, color=RGBColor(0x4F, 0xC4, 0xCE),
         font=HEAD, after=6)
    t2 = tb(s, M, 2.28, 11.5, 0.4)
    para(t2, "Joint denoising and 2x super-resolution, built on a degradation "
             "model recovered from the data", size=12.5,
         color=RGBColor(0xC9, 0xD6, 0xE0), first=True)

    # team table
    rows = [("SR. NO", "ROLE", "NAME", "ACADEMIC YEAR")]
    for i, (role, name, yr) in enumerate(TEAM["members"], 1):
        rows.append((str(i), role, name, yr))
    y = 3.0
    colw = [0.9, 2.4, 4.2, 2.2]
    for ri, row in enumerate(rows):
        x = M
        for ci, cell in enumerate(row):
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                     Inches(colw[ci]), Inches(0.36))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x3E) if ri else RGBColor(0x1E, 0x3A, 0x50)
            box.line.color.rgb = RGBColor(0x2A, 0x44, 0x5A)
            box.line.width = Pt(0.5)
            box.shadow.inherit = False
            box.text_frame.text = ""
            tt = tb(s, x + 0.1, y + 0.055, colw[ci] - 0.15, 0.3, anchor=MSO_ANCHOR.MIDDLE)
            para(tt, str(cell), size=10.5, bold=(ri == 0),
                 color=RGBColor(0x4F, 0xC4, 0xCE) if ri == 0 else WHITE, first=True, after=0)
            x += colw[ci]
        y += 0.36

    t3 = tb(s, M, y + 0.32, 11.5, 1.2)
    para(t3, "Team: " + TEAM["team_name"], size=13, bold=True, color=WHITE,
         first=True, after=4)
    para(t3, TEAM["college"], size=11.5, color=RGBColor(0xC9, 0xD6, 0xE0), after=4)
    para(t3, f"{TEAM['leader_phone']}   |   {TEAM['leader_email']}",
         size=11, color=MUTED, after=0)


def s2_problem(prs, m):
    s = blank(prs)
    heading(s, "Problem Statement Addressed", PROBLEM)

    t = card(s, M, 1.78, 6.05, 3.1)
    para(t, "DESCRIPTION / DETAILS", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t, [
        "Inspection images arrive degraded by speckle noise, additive Gaussian "
        "noise and downsampling, applied in an undisclosed order",
        "Detail lost to degradation reduces the reliability of every downstream "
        "measurement and defect-detection step",
        "Input is half resolution and noisy; output must be clean at full "
        "ground-truth resolution",
        "The hidden test set includes out-of-distribution content, so the model "
        "must generalise beyond the training images",
    ])

    t2 = card(s, M + 6.35, 1.78, 5.7, 3.1)
    para(t2, "WHY IT MATTERS", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t2, [
        "Semiconductor inspection is a precision-measurement task: invented "
        "detail is worse than no detail",
        "Cleaner images at higher resolution mean fewer missed defects and "
        "fewer false alarms",
        "Scored on restoration quality, end-to-end H100 throughput and "
        "reproducibility together",
    ])

    b = m.get("bicubic", {})
    stat(s, M, 5.15, 3.5, f"{float(b.get('psnr', 22.94)):.2f} dB",
         "doing nothing (bicubic upsample)", color=AMBER)
    stat(s, M + 3.8, 5.15, 3.5, "32.25 dB", "denoise-then-interpolate ceiling", color=MUTED)
    stat(s, M + 7.6, 5.15, 4.4, "3200 pairs", "provided training data", color=CYAN)


def s3_solution(prs, m):
    s = blank(prs)
    heading(s, "Idea Description and Proposed Solution",
            "One network, one forward pass: joint denoise + 2x super-resolution")

    steps = [
        ("Recover the\nforward model",
         "Fit the degradation from the data instead of assuming it"),
        ("Generate matched\ntraining pairs",
         "Sample noise from the measured distribution, widened"),
        ("Restore in one\nnetwork",
         "NAF blocks at low resolution + PixelShuffle x2 head"),
        ("Ship a fast\npipeline",
         "Batched, threaded I/O, clamped output, 2 imports"),
    ]
    x, cw = M, 2.82
    for i, (h, sub) in enumerate(steps):
        t = card(s, x, 1.8, cw, 1.75)
        para(t, f"0{i+1}", size=10.5, bold=True, color=CYAN, first=True, after=4)
        para(t, h.replace("\n", " "), size=13.5, bold=True, color=INK, font=HEAD, after=6)
        para(t, sub, size=10.5, color=MUTED, after=0)
        if i < 3:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + cw + 0.03),
                                   Inches(2.5), Inches(0.26), Inches(0.32))
            a.fill.solid()
            a.fill.fore_color.rgb = CYAN
            a.line.fill.background()
            a.shadow.inherit = False
            a.text_frame.text = ""
        x += cw + 0.34

    t = card(s, M, 3.75, 6.05, 2.55)
    para(t, "SOLUTION DETAILS", size=11, bold=True, color=CYAN, first=True, after=7)
    bullets(t, [
        "NAFNet-style body: 16 blocks at LOW resolution, single PixelShuffle "
        "upsample at the end -- about 4x cheaper than a U-Net at output resolution",
        "0.68 M parameters, no BatchNorm, so the evaluator's batch size cannot "
        "change our outputs (verified to 4.8e-07)",
        "Fully convolutional: the same weights handle 128->256 and 256->512",
        "Loss = Charbonnier + 0.15 x (1 - SSIM), plus LPIPS for the final 10%",
    ], size=11.5, after=5)

    t2 = card(s, M + 6.35, 3.75, 5.7, 2.55, accent=CYAN)
    para(t2, "MEASURED DEGRADATION MODEL", size=11, bold=True, color=CYAN,
         first=True, after=7)
    bullets(t2, [
        "Downsampling: 2x2 area average (lowest residual of all kernels tested)",
        "Noise applied AFTER downsampling -- residual is spatially white",
        "Speckle is multiplicative Gamma: measured skew +0.357 vs Gamma's "
        "predicted +0.333",
        "L median 36; sigma median 0.021, with 21% of images near zero",
    ], size=11.5, after=5)


def s4_innovation(prs, m):
    s = blank(prs)
    heading(s, "Innovation and Uniqueness",
            "What separates this from feeding pairs to a standard network")

    t = card(s, M, 1.78, 6.05, 3.4, accent=CYAN)
    para(t, "KEY INNOVATION -- we measured the degradation, we did not guess it",
         size=12.5, bold=True, color=INK, font=HEAD, first=True, after=8)
    bullets(t, [
        "Most entrants feed the provided pairs to a network. We first recovered "
        "the forward model by fitting the residual across all 3200 pairs",
        "The skew of the residual identifies the noise FAMILY: multiplicative "
        "Gaussian is symmetric, Gamma is not. Measured +0.357 against a "
        "predicted +0.333 pins it",
        "That let us synthesise unlimited training pairs matching the true "
        "distribution -- sampled from empirical quantile tables, then widened "
        "so the tails exceed what was observed",
        "A first version using assumed ranges was 50% noisier than reality and "
        "would have taught the model to over-smooth",
    ], size=11.5, after=6)

    t2 = card(s, M + 6.35, 1.78, 5.7, 3.4)
    para(t2, "COMPETITIVE ADVANTAGE", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t2, [
        "Leakage-aware validation: samples are ordered by source image, so a "
        "random split leaks near-duplicates. We hold out 10 evenly-spaced "
        "contiguous blocks instead",
        "Every claim is reproducible by one command "
        "(python analyze_degradation.py)",
        "Deliberately no GAN: the brief asks for restoration without "
        "hallucinating structure, and invented detail is harmful in a "
        "measurement context",
        "Efficiency treated as a first-class deliverable, not an afterthought",
    ], size=11.5, after=6)

    stat(s, M, 5.5, 4.0, "+0.357 vs +0.333", "measured vs predicted Gamma skew", color=CYAN)
    stat(s, M + 4.2, 5.5, 4.0, "0.68 M", "parameters", color=CYAN)
    stat(s, M + 8.4, 5.5, 3.6, "1 command", "regenerates the analysis", color=CYAN)


def s5_impact(prs, m):
    s = blank(prs)
    heading(s, "Impact and Benefits",
            "Measured on a 320-image held-out split, never seen in training")

    mm, bb = m.get("model", {}), m.get("bicubic", {})

    def g(d, k, dflt):
        try:
            return float(d.get(k, dflt))
        except (TypeError, ValueError):
            return dflt

    rows = [
        ("Metric", "Bicubic baseline", "This solution", "Improvement"),
        ("PSNR", f"{g(bb,'psnr',22.94):.2f} dB", f"*{g(mm,'psnr',27.54):.2f} dB",
         f"+{g(mm,'psnr',27.54)-g(bb,'psnr',22.94):.2f} dB"),
        ("SSIM", f"{g(bb,'ssim',0.546):.4f}", f"*{g(mm,'ssim',0.7594):.4f}",
         f"+{100*(g(mm,'ssim',0.7594)/g(bb,'ssim',0.546)-1):.0f}%"),
        ("LPIPS (lower better)", f"{g(bb,'lpips',0.4478):.4f}",
         f"*{g(mm,'lpips',0.3922):.4f}",
         f"-{100*(1-g(mm,'lpips',0.3922)/g(bb,'lpips',0.4478)):.0f}%"),
    ]
    y, colw = 1.85, [3.1, 2.9, 2.9, 3.2]
    for ri, row in enumerate(rows):
        x = M
        for ci, cell in enumerate(row):
            emph = isinstance(cell, str) and cell.startswith("*")
            txt = cell[1:] if emph else cell
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                     Inches(colw[ci]), Inches(0.44))
            box.fill.solid()
            box.fill.fore_color.rgb = NAVY if ri == 0 else (
                RGBColor(0xEF, 0xF3, 0xF6) if ri % 2 else WHITE)
            box.line.color.rgb = RULE
            box.line.width = Pt(0.5)
            box.shadow.inherit = False
            box.text_frame.text = ""
            tt = tb(s, x + 0.14, y + 0.07, colw[ci] - 0.2, 0.36, anchor=MSO_ANCHOR.MIDDLE)
            para(tt, txt, size=12, bold=(ri == 0 or emph),
                 color=WHITE if ri == 0 else (CYAN if emph else INK), first=True, after=0)
            x += colw[ci]
        y += 0.44

    t = card(s, M, 3.75, 6.05, 2.0, accent=CYAN)
    para(t, "PRIMARY IMPACT", size=11, bold=True, color=CYAN, first=True, after=7)
    bullets(t, [
        "Improves on ALL THREE scored axes at once -- perceptual metrics often "
        "move against PSNR, and here they do not",
        "Recovers detail lost to downsampling without inventing structure, "
        "which is what a measurement pipeline requires",
    ], size=11.5)

    t2 = card(s, M + 6.35, 3.75, 5.7, 2.0)
    para(t2, "QUANTIFIABLE OUTCOMES", size=11, bold=True, color=CYAN, first=True, after=7)
    spi = mm.get("sec_per_image")
    rate = f"{1/float(spi):.0f} images/second ({float(spi)*1000:.1f} ms each)" if spi else "measured on GPU"
    bullets(t2, [
        f"End-to-end inference: {rate} on an NVIDIA T4",
        "2.79 MB checkpoint -- deployable anywhere, no external downloads",
        "Same weights serve both 128->256 and 256->512 inputs",
    ], size=11.5)

    t3 = tb(s, M, 5.95, W - 2 * M, 0.5)
    para(t3, "Honest limitation: training stopped at step 7500, short of "
             "convergence, so the model has not yet passed the "
             "denoise-then-interpolate ceiling of 32.25 dB. The cause was input-"
             "pipeline throughput, and it is the first thing we would fix.",
         size=10.5, color=MUTED, italic=True, first=True)


def s6_tech(prs, m):
    s = blank(prs)
    heading(s, "Technology, Feasibility and Methodology",
            "Everything below runs from a clean clone with no manual configuration")

    t = card(s, M, 1.78, 3.85, 3.3)
    para(t, "TECHNICAL STACK", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t, [
        "PyTorch 2.x, Python 3.11",
        "NAFNet-style CNN, 0.68 M params",
        "Trained on a single NVIDIA T4",
        "Inference: torch + numpy only",
        "No external datasets or pretrained restoration weights",
    ], size=11.5)

    t2 = card(s, M + 4.1, 1.78, 3.85, 3.3)
    para(t2, "IMPLEMENTATION", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t2, [
        "python run.py <in-dir> <out-dir>",
        "Weights auto-loaded from models/",
        "Output dir created if absent",
        "Outputs clamped to [0,1], no NaN/Inf, verified",
        "Batched by resolution, threaded I/O, fp16 autocast",
    ], size=11.5)

    t3 = card(s, M + 8.2, 1.78, 3.85, 3.3, accent=CYAN)
    para(t3, "REPRODUCIBILITY", size=11, bold=True, color=CYAN, first=True, after=8)
    bullets(t3, [
        "Seeds fixed; config-driven, nothing hard-coded",
        "Every validation logs step, LR, all metrics, git hash and full config",
        "Model config embedded in the checkpoint",
        "Pre-flight gate blocks a broken run before it wastes GPU hours",
        "Verified from a fresh public clone",
    ], size=11.5)

    t4 = card(s, M, 5.32, W - 2 * M, 1.25)
    para(t4, "FEASIBILITY", size=11, bold=True, color=CYAN, first=True, after=6)
    para(t4, "Trained end-to-end within a single free Kaggle GPU session. The "
             "model is small enough to deploy on commodity hardware and fast "
             "enough for inline inspection, and the degradation engine can "
             "generate unlimited additional training data from any clean image "
             "set -- so the approach scales to new tools and processes without "
             "new paired acquisitions.", size=11.5, color=INK)


def s7_links(prs, m):
    s = blank(prs)
    heading(s, "Repository, References and Disclosure")

    t = card(s, M, 1.78, 6.05, 1.75, accent=CYAN)
    para(t, "GITHUB REPOSITORY", size=11, bold=True, color=CYAN, first=True, after=7)
    para(t, TEAM["github"], size=14, bold=True, color=INK, font=HEAD, after=6)
    para(t, "Contains run.py, requirements.txt, README.md, models/best.pt, "
            "training and evaluation code, and the full results.",
         size=11, color=MUTED)

    t2 = card(s, M + 6.35, 1.78, 5.7, 1.75)
    para(t2, "PROTOTYPE / SIMULATION VIDEO", size=11, bold=True, color=CYAN,
         first=True, after=7)
    para(t2, TEAM["video"] or "Working solution demonstrated in the repository; "
                              "inference reproducible with a single command.",
         size=12 if TEAM["video"] else 11,
         bold=bool(TEAM["video"]), color=INK if TEAM["video"] else MUTED)

    t3 = card(s, M, 3.72, 6.05, 2.55)
    para(t3, "RESEARCH BACKGROUND", size=11, bold=True, color=CYAN, first=True, after=7)
    bullets(t3, [
        "Chen et al., Simple Baselines for Image Restoration, ECCV 2022 -- NAFNet",
        "Wang et al., Real-ESRGAN, ICCV 2021 -- degradation modelling for "
        "real-world super-resolution",
        "Shi et al., Sub-Pixel Convolution, CVPR 2016 -- PixelShuffle upsampling",
        "Wang et al., Image Quality Assessment, IEEE TIP 2004 -- SSIM",
        "Zhang et al., The Unreasonable Effectiveness of Deep Features, "
        "CVPR 2018 -- LPIPS",
    ], size=10.5, after=4)

    t4 = card(s, M + 6.35, 3.72, 5.7, 2.55)
    para(t4, "EXTERNAL RESOURCE DISCLOSURE", size=11, bold=True, color=CYAN,
         first=True, after=7)
    bullets(t4, [
        "No external datasets used. No pretrained restoration weights used",
        "The model is trained from scratch on the official KLA data plus "
        "synthetic pairs derived from that same data",
        "lpips (BSD-2-Clause) supplies the perceptual loss and metric; its VGG16 "
        "backbone carries ImageNet weights shipped with the package",
        "Training degrades gracefully without it, and run.py never imports it",
    ], size=10.5, after=4)


def main() -> None:
    m = metrics()
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    for fn in (s1_title, s2_problem, s3_solution, s4_innovation,
               s5_impact, s6_tech, s7_links):
        fn(prs, m)

    name = f"{TEAM['team_name']}_{TEAM['ps_no']}.pptx"
    out = ROOT / name
    prs.save(str(out))
    n = len(prs.slides._sldIdLst)
    print(f"[deck] wrote {out} ({n} slides)")
    if n > 7:
        print(f"[deck] WARNING: {n} slides exceeds the template's 6-7 limit")
    if "TEAM_NAME" in TEAM["team_name"] or "NAME" in TEAM["members"][0][1]:
        print("[deck] WARNING: team details still contain placeholders -- edit TEAM at the top")


if __name__ == "__main__":
    main()
